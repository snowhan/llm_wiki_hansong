#!/usr/bin/env python3
"""
extract_images.py — Extract embedded images from PDF, DOCX, and PPTX files.

Usage:
    python3 extract_images.py <input_file> <output_dir>

Output:
    JSON to stdout: {"images": ["/abs/path/to/img0.png", ...], "error": null}
    On failure:     {"images": [], "error": "<message>"}

Dependencies (install with pip):
    pymupdf   — PDF support (pip install pymupdf)
    python-docx — DOCX support (pip install python-docx)
    python-pptx — PPTX support (pip install python-pptx)
    Pillow    — image write fallback (pip install Pillow)

Images are saved as PNG files named: 0.png, 1.png, 2.png, ...
A maximum of MAX_IMAGES images are extracted to control token costs.
"""

import sys
import os
import json
import hashlib

MAX_IMAGES = 10


def extract_from_pdf(file_path: str, output_dir: str) -> list[str]:
    try:
        import fitz  # PyMuPDF
    except ImportError:
        raise RuntimeError("pymupdf is not installed. Run: pip install pymupdf")

    doc = fitz.open(file_path)
    saved: list[str] = []

    for page_num in range(len(doc)):
        if len(saved) >= MAX_IMAGES:
            break
        page = doc[page_num]
        image_list = page.get_images(full=True)
        for img_index, img in enumerate(image_list):
            if len(saved) >= MAX_IMAGES:
                break
            xref = img[0]
            base_image = doc.extract_image(xref)
            image_bytes = base_image["image"]
            ext = base_image.get("ext", "png")
            img_path = os.path.join(output_dir, f"{len(saved)}.{ext}")
            with open(img_path, "wb") as f:
                f.write(image_bytes)
            saved.append(img_path)

    doc.close()
    return saved


def extract_from_docx(file_path: str, output_dir: str) -> list[str]:
    try:
        from docx import Document
    except ImportError:
        raise RuntimeError("python-docx is not installed. Run: pip install python-docx")

    doc = Document(file_path)
    saved: list[str] = []

    for rel in doc.part.rels.values():
        if len(saved) >= MAX_IMAGES:
            break
        if "image" in rel.reltype:
            img_part = rel.target_part
            ext = img_part.content_type.split("/")[-1]
            if ext == "jpeg":
                ext = "jpg"
            img_path = os.path.join(output_dir, f"{len(saved)}.{ext}")
            with open(img_path, "wb") as f:
                f.write(img_part.blob)
            saved.append(img_path)

    return saved


def extract_from_pptx(file_path: str, output_dir: str) -> list[str]:
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError:
        raise RuntimeError("python-pptx is not installed. Run: pip install python-pptx")

    prs = Presentation(file_path)
    saved: list[str] = []

    for slide in prs.slides:
        if len(saved) >= MAX_IMAGES:
            break
        for shape in slide.shapes:
            if len(saved) >= MAX_IMAGES:
                break
            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                img = shape.image
                ext = img.ext
                img_path = os.path.join(output_dir, f"{len(saved)}.{ext}")
                with open(img_path, "wb") as f:
                    f.write(img.blob)
                saved.append(img_path)

    return saved


def main():
    if len(sys.argv) < 3:
        print(json.dumps({"images": [], "error": "Usage: extract_images.py <input_file> <output_dir>"}))
        sys.exit(1)

    file_path = sys.argv[1]
    output_dir = sys.argv[2]

    if not os.path.exists(file_path):
        print(json.dumps({"images": [], "error": f"File not found: {file_path}"}))
        sys.exit(1)

    os.makedirs(output_dir, exist_ok=True)

    ext = os.path.splitext(file_path)[1].lower()

    try:
        if ext == ".pdf":
            images = extract_from_pdf(file_path, output_dir)
        elif ext in (".docx", ".doc"):
            images = extract_from_docx(file_path, output_dir)
        elif ext in (".pptx", ".ppt"):
            images = extract_from_pptx(file_path, output_dir)
        else:
            images = []

        print(json.dumps({"images": images, "error": None}))

    except Exception as e:
        print(json.dumps({"images": [], "error": str(e)}))
        sys.exit(1)


if __name__ == "__main__":
    main()
